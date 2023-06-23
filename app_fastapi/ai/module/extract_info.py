from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# from pdfminer.high_level import extract_text, extract_pages
# from pdfminer.layout import LTTextContainer
# import PyPDF2
# import fitz
import fitz
from PIL import Image

from collections import defaultdict
from kiwipiepy import Kiwi

import re
import os
import shutil
from bs4 import BeautifulSoup

from ai.module.util import FileUtil


class ExtractInfo:
    def __init__(self) -> None:
        # self.tp = Text_Preprocessing()
        self.fileutil = FileUtil()
        self.kiwi = Kiwi()
        self.reset()

    def reset(self):
        self.file_name = None
        self.name = None
        self.ext = None
        self.text_dict = {'total': [], 'page': defaultdict(list)}
        self.link_list = []
        self.img_dict = defaultdict(list)
        self.text_info = None
        self.page_num = None
        self.img_num = 0
        self.image_blob = []
        self.pre_text = None


class PPT_Info_Extract(ExtractInfo):
    def __init__(self) -> None:
        super().__init__()

    def reset(self):
        super().reset()
        self.parsed = None

    def run(self, file_name):
        """_summary_

        Args:
            file_name (str): 문서 이름

        Returns:
            dict:
                page_num (int):
                    text (str): list
                    img (str): list
        """
        self._extract(file_name)
        res = {}
        texts = self.text_dict['page']
        imgs = self.img_dict

        pages = set(list(texts.keys()) + list(imgs.keys()))
        for page in range(min(pages), max(pages)+1):
            res[page] = {}
            res[page]['text'] = texts[page]
            res[page]['img'] = imgs[page]
        return res

    def _extract(self, file_name):
        # path = "deep_learning_intro.pptx"
        self.reset()
        self.file_name = file_name
        self.name, self.ext = os.path.splitext(file_name)
        path = self.fileutil.orignal_dir + file_name
        self.parsed = Presentation(path)

        for idx, slide in enumerate(self.parsed.slides, start=1):
            self.page_num = idx
            for shape in slide.shapes:
                self._group_check(shape)

    def _group_check(self, shape):
        # 그룹 모형
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                self._info_extract(s)
        # 그룹 모형 아님
        else:
            self._info_extract(shape)

    def _info_extract(self, shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            self._group_check(shape)
        else:
            # yes image
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = shape.image.blob

                    # 이미지 중복 체크
                    if image_blob not in self.image_blob:
                        ext = shape.image.ext  # - (확장명)
                        size = shape.image.size
                        if ext in ['png', 'jpeg', 'jpg'] and size[0] >= 200 and size[1] >= 200:
                            self.img_num += 1
                            num = str(self.img_num).zfill(3)
                            save_path = f"{self.fileutil.res_dir}{self.name}/image_{num}.{ext}"
                            with open(save_path, "wb") as file:
                                file.write(image_blob)
                            self.image_blob.append(image_blob)
                            self.img_dict[self.img_num].append(save_path)
                except:
                    pass
            # no image
            else:
                # yes text
                if shape.has_text_frame:
                    if shape.text.strip() != "":
                        self.text_dict['page'][self.page_num].append(
                            shape.text.strip())


class DOCX_Info_Extract(ExtractInfo):
    def __init__(self) -> None:
        super().__init__()

    def reset(self):
        super().reset()
        self.doc = None

    def run(self, file_name):
        """_summary_

        Args:
            file_name (str): 문서 이름

        Returns:
            dict:
                page_num (int):
                    text (str): list
                    img (str): list
        """
        self._extract(file_name)
        res = {}
        texts = self.text_dict['page']
        imgs = self.img_dict

        self.text_dict['page'][self.page_num]
        self.text_dict['page'][self.page_num].append(
            self._extract_tables())
        self._extract_images()

        pages = set(list(texts.keys()) + list(imgs.keys()))
        for page in range(min(pages), max(pages)+1):
            res[page] = {}
            res[page]['text'] = texts[page]
            res[page]['img'] = imgs[page]
        return res

    def _extract(self, file_name):
        # path = "deep_learning_intro.pptx"
        self.reset()
        self.file_name = file_name
        self.name, self.ext = os.path.splitext(file_name)
        path = self.fileutil.orignal_dir + file_name
        self.doc = Document(path)

        self.text_dict['page'][self.page_num]
        self.text_dict['page'][self.page_num].append(
            self._extract_tables())
        self._extract_images()

    def _extract_text(self):
        paragraphs = [p.text for p in self.doc.paragraphs]
        return '\n'.join(paragraphs)

    def _extract_tables(self):
        tables = []
        for table in self.doc.tables:
            data = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text)
                data.append(' '.join(cells))
            tables.append(' '.join(data))
        return '\n'.join(tables)

    def _extract_images(self):
        for rel in self.doc.part.rels.values():
            if "image" in rel.reltype:
                image = rel.target_part.image
                image_blob = image.blob

                # 이미지 중복 체크
                if image_blob not in self.image_blob:
                    ext = image.ext
                    size = image.px_height, image.px_width
                    if ext in ['png', 'jpeg', 'jpg'] and size[0] >= 200 and size[1] >= 200:
                        self.img_num += 1
                        num = str(self.img_num).zfill(3)
                        save_path = f"{self.fileutil.res_dir}{self.name}/image_{num}.{ext}"

                        with open(save_path, "wb") as file:
                            file.write(image_blob)
                        self.image_blob.append(image_blob)
                        self.img_dict[self.img_num].append(save_path)


class PDF_Info_Extract(ExtractInfo):
    def __init__(self) -> None:
        super().__init__()

    def reset(self):
        super().reset()
        self.pdf = None

    def run(self, file_name):
        """_summary_

        Args:
            file_name (str): 문서 이름

        Returns:
            dict:
                page_num (int):
                    text (str): list
                    img (str): list
        """
        self._extract(file_name)
        res = {}
        texts = self.text_dict['page']
        imgs = self.img_dict

        pages = set(list(texts.keys()) + list(imgs.keys()))
        for page in range(min(pages), max(pages)+1):
            res[page] = {}
            res[page]['text'] = texts[page]
            res[page]['img'] = imgs[page]
        return res

    def _extract(self, file_name):
        # path = "deep_learning_intro.pptx"
        self.reset()
        self.file_name = file_name
        self.name, self.ext = os.path.splitext(file_name)
        path = self.fileutil.orignal_dir + file_name
        self.pdf = fitz.open(path)

        for idx, page_content in enumerate(self.pdf, start=1):
            self.page_num = idx
            # 텍스트 추출
            self._extract_text(page_content)
            # 이미지 추출
            self._extract_image(page_content)

    def _extract_text(self, page_content):
        text = page_content.get_text()
        if text != "":
            self.text_dict['page'][self.page_num].append(
                text.strip())

    def _extract_image(self, page_content):
        images_list = page_content.get_images()
        for img_info in images_list:
            # Extract image
            image = self.pdf.extract_image(img_info[0])
            ext = image['ext']
            size = image['height'], image['width']
            if ext in ['png', 'jpeg', 'jpg'] and size[0] >= 200 and size[1] >= 200:
                self.img_num += 1
                num = str(self.img_num).zfill(3)
                save_path = f"{self.fileutil.res_dir}{self.name}/image_{num}.{ext}"

                # Store image bytes
                image_blob = image['image']

                # Save image
                with open(save_path, 'wb') as file:
                    file.write(image_blob)
                self.image_blob.append(image_blob)
                self.img_dict[self.img_num].append(save_path)


if __name__ == '__main__':
    # extractinfo = ExtractInfo()

    # path = ""

    # res = extractinfo.run(path)
    # print(os.listdir('db/data/input/'))
    fileutil = FileUtil()
    fileutil.input_files_update()

    # ppt = PPT_Info_Extract()
    # ppt.extract('1 Intro.pptx')
    # ppt = DOCX_Info_Extract()
    # ppt.extract('판넬양식.docx')
    # print(dict(ppt.text_dict['page']))
    # print(dict(ppt.img_dict))
    # print(ppt.run()[10])
